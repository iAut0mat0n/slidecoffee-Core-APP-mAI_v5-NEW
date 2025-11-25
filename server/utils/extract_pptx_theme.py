#!/usr/bin/env python3
"""
PowerPoint Theme Extractor

Extracts theme information (colors, fonts, images) from .pptx files
using python-pptx library.

Usage:
    python extract_pptx_theme.py <path_to_pptx_file>
    
Output:
    JSON with extracted theme data
"""

import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.xmlchemy import OxmlElement
import os
import base64

def rgb_to_hex(rgb_tuple):
    if not rgb_tuple or len(rgb_tuple) < 3:
        return None
    return '#{:02x}{:02x}{:02x}'.format(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])

def extract_colors_from_slide_master(prs):
    colors = {
        'primary': None,
        'secondary': None,
        'accent': None,
        'background': None,
        'text': None,
        'palette': []
    }
    
    try:
        if prs.slide_master:
            slide_master = prs.slide_master
            
            for shape in slide_master.shapes:
                if hasattr(shape, 'fill'):
                    if shape.fill.type == 1:
                        try:
                            rgb = shape.fill.fore_color.rgb
                            hex_color = rgb_to_hex((rgb[0], rgb[1], rgb[2]))
                            if hex_color and hex_color not in colors['palette']:
                                colors['palette'].append(hex_color)
                                if not colors['primary']:
                                    colors['primary'] = hex_color
                                elif not colors['secondary']:
                                    colors['secondary'] = hex_color
                                elif not colors['accent']:
                                    colors['accent'] = hex_color
                        except:
                            pass
    except:
        pass
    
    if not colors['primary']:
        colors['primary'] = '#7C3AED'
    if not colors['secondary']:
        colors['secondary'] = '#6EE7B7'
    if not colors['accent']:
        colors['accent'] = '#FFE5E5'
    
    return colors

def extract_fonts(prs):
    fonts = {
        'heading': None,
        'body': None,
        'fonts_used': set()
    }
    
    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                fonts['fonts_used'].add(run.font.name)
                                if not fonts['heading'] and paragraph.level == 0:
                                    fonts['heading'] = run.font.name
                                elif not fonts['body']:
                                    fonts['body'] = run.font.name
    except:
        pass
    
    fonts['fonts_used'] = list(fonts['fonts_used'])
    
    if not fonts['heading']:
        fonts['heading'] = fonts['body'] if fonts['body'] else 'Inter'
    if not fonts['body']:
        fonts['body'] = 'Inter'
    
    return fonts

def extract_images(prs):
    images = []
    logo_image = None
    
    try:
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        if len(image_bytes) < 50000 and slide_idx == 0:
                            logo_image = {
                                'slide_number': slide_idx + 1,
                                'shape_index': shape_idx,
                                'content_type': image.content_type,
                                'size_bytes': len(image_bytes),
                                'potential_logo': True
                            }
                        
                        images.append({
                            'slide_number': slide_idx + 1,
                            'shape_index': shape_idx,
                            'content_type': image.content_type,
                            'size_bytes': len(image_bytes)
                        })
                    except:
                        pass
    except:
        pass
    
    return {
        'total_images': len(images),
        'images': images[:5],
        'logo_image': logo_image
    }

def extract_slide_layouts(prs):
    layouts = []
    
    try:
        if prs.slide_master:
            for layout in prs.slide_master.slide_layouts:
                layouts.append({
                    'name': layout.name,
                    'shape_count': len(layout.shapes)
                })
    except:
        pass
    
    return layouts

def extract_theme(pptx_path):
    try:
        prs = Presentation(pptx_path)
        
        colors = extract_colors_from_slide_master(prs)
        fonts = extract_fonts(prs)
        images = extract_images(prs)
        layouts = extract_slide_layouts(prs)
        
        slide_dimensions = {
            'width': prs.slide_width,
            'height': prs.slide_height,
            'width_inches': prs.slide_width / Inches(1),
            'height_inches': prs.slide_height / Inches(1)
        }
        
        theme_data = {
            'success': True,
            'filename': os.path.basename(pptx_path),
            'colors': colors,
            'fonts': fonts,
            'images': images,
            'slide_layouts': layouts,
            'slide_dimensions': slide_dimensions,
            'total_slides': len(prs.slides),
            'has_master_slides': bool(prs.slide_master),
            'metadata': {
                'core_properties': {
                    'title': prs.core_properties.title if prs.core_properties.title else None,
                    'author': prs.core_properties.author if prs.core_properties.author else None,
                    'subject': prs.core_properties.subject if prs.core_properties.subject else None
                }
            }
        }
        
        return theme_data
        
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'error_type': type(e).__name__
        }

def main():
    if len(sys.argv) < 2:
        print(json.dumps({
            'success': False,
            'error': 'No file path provided. Usage: python extract_pptx_theme.py <path_to_pptx>'
        }))
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    
    if not os.path.exists(pptx_path):
        print(json.dumps({
            'success': False,
            'error': f'File not found: {pptx_path}'
        }))
        sys.exit(1)
    
    result = extract_theme(pptx_path)
    print(json.dumps(result, indent=2))

if __name__ == '__main__':
    main()
